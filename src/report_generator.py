from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from datetime import datetime
import pandas as pd
import matplotlib.pyplot as plt
import os

def create_ppt_report(data):
    """금융 데이터를 기반으로 PPT 보고서를 생성하는 함수"""
    file_path = "output/daily_summary.pptx"
    if os.path.exists(file_path):
        os.remove(file_path)
        print("📄 기존 PPT 파일 삭제 완료!")
    
    prs = Presentation()

    layout=prs.slide_layouts[6] # 빈 슬라이드 레이아웃
    slide=prs.slides.add_slide(prs.slide_layouts[6]) # 빈 슬라이드 레이아웃을 가진 슬라이드 1장 추가
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 245) # 배경색 변경

    # A4 크기로 조정
    prs.slide_width=Cm(21.6)
    prs.slide_height=Cm(27)

    # 제목
    textbox = slide.shapes.add_textbox(Cm(1.8), Cm(0.85), Cm(7.15), Cm(1.9))
    text_frame = textbox.text_frame
    p = text_frame.add_paragraph()
    run = p.add_run()
    run.text = "오늘의 주요 지표"
    run.font.name = "Segoe UI"
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = RGBColor(30, 58, 95)
    p = text_frame.add_paragraph()

    run = p.add_run()
    today = datetime.today()
    formatted_date = today.strftime("%Y년 %m월 %d일 %A")
    run.text = formatted_date
    run.font.name = "Segoe UI"
    run.font.size = Pt(18)
    run.font.bold = False
    run.font.color.rgb = RGBColor(30, 58, 95)

    name_list = ['S&P500', 'NASDAQ', 'Dow Jones', 'US 10Y Yield', 'USD/KRW', 'KOSPI']

    #종목별 정보 추가
    for stock in data.iterrows():
        y_dir = 6
        if(stock[0] // 3 == 1):
            x_dir = 13.3
        else:
            x_dir = 3
        y_dir += stock[0] * 6 - stock[0] // 3 * 18
        textbox = slide.shapes.add_textbox(Cm(x_dir), Cm(y_dir), Cm(4.85), Cm(1.45))
        text_frame = textbox.text_frame
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE # 수직 가운데 정렬
        p = text_frame.add_paragraph()
        p.alignment = PP_ALIGN.CENTER # 수평 가운데 정렬
        
        # 종목명
        run = p.add_run()
        run.text = name_list[stock[0]]
        run.font.name = "Segoe UI"
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.color.rgb = RGBColor(30, 58, 95)

        # 종가
        run = p.add_run()
        if(stock[0] == 3):
            run.text = f"\n{round(stock[1]['최근 종가'], 3)}" # 10년물은 소수점 세자리리
        else:
            run.text = f"\n{round(stock[1]['최근 종가'], 2)}"
        run.font.name = "Segoe UI"
        run.font.size = Pt(48)
        run.font.bold = True
        if stock[1]['변동'] > 0:
            run.font.color.rgb = RGBColor(0, 200, 0)
        else:
            run.font.color.rgb = RGBColor(255, 0, 0)

        # 변동량과 변동률
        run = p.add_run()
        if(stock[0] == 3):
            run.text = f"\n{round(stock[1]['변동'], 3):+} ({round(stock[1]['변동률'], 2):+}%)" # 10년물은 소수점 세자리
        else:
            run.text = f"\n{round(stock[1]['변동'], 2):+} ({round(stock[1]['변동률'], 2):+}%)"
        run.font.name = "Segoe UI"
        run.font.size = Pt(20)
        if stock[1]['변동'] > 0:
            run.font.color.rgb = RGBColor(0, 200, 0)
        else:
            run.font.color.rgb = RGBColor(255, 0, 0)
    
    prs.save(file_path)
    print("📄 PPT 파일 생성 완료!")

if __name__ == "__main__":
    df = pd.read_csv("data/stock_data.csv", encoding='utf-8-sig')
    create_ppt_report(df)